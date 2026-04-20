# from pptx2md import convert, ConversionConfig
# from pathlib import Path

# # Basic usage
# convert(
#     ConversionConfig(
#         pptx_path=Path('/home/chenyifan/rag/ppt_to_md/2025_11.20陈奕帆组会.pptx'),
#         output_path=Path('output.md'),
#         image_dir=Path('img'),
#         disable_notes=True
#     )
# )

from pathlib import Path
from pptx import Presentation
import subprocess
import re


# def ppt_to_md(ppt_path: str | Path) -> Path:
#     ppt_path = Path(ppt_path).resolve()

#     output_path = ppt_path.with_suffix(".md").resolve()
#     image_dir = (ppt_path.parent / (ppt_path.stem + "_img")).resolve()

#     image_dir.mkdir(parents=True, exist_ok=True)

#     config = ConversionConfig(
#         pptx_path=ppt_path,
#         output_path=output_path,
#         image_dir=image_dir,
#         disable_notes=True,
#     )

#     convert(config)

#     return output_path



def ppt_to_md(ppt_path: str | Path) -> Path:
    ppt_path = Path(ppt_path).resolve()
    output_path = ppt_path.with_suffix(".md").resolve()

    prs = Presentation(ppt_path)
    md_lines = []

    for i, slide in enumerate(prs.slides, start=1):
        # 幻灯片标题
        title_shapes = [sh for sh in slide.shapes if sh.has_text_frame and sh.text.strip()]
        if title_shapes:
            slide_title = title_shapes[0].text.strip()
            md_lines.append(f"# {slide_title}\n")
        
        # 幻灯片正文
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue  # 忽略图片
            text = shape.text.strip()
            if text and text != slide_title:  # 避免重复写标题
                md_lines.append(text + "\n")
        
        # 幻灯片分页符
        md_lines.append("\n---\n")

    # 写入文件
    output_path.write_text("\n".join(md_lines), encoding="utf-8")
    return output_path

def pdf_to_md(ppt_path: str | Path) -> Path:

    pdf_path = Path(ppt_path).resolve()
    output_dir = pdf_path.parent

    # 1. 运行 marker
    subprocess.run(
        [
            "marker_single",
            str(pdf_path),
            "--output_dir",
            str(output_dir),
        ],
        check=True,
    )

    # marker生成的目录
    md_dir = output_dir / pdf_path.stem
    md_path = md_dir / f"{pdf_path.stem}.md"

    if not md_path.exists():
        raise FileNotFoundError(f"Markdown file not found: {md_path}")

    # 2. 清理 markdown 内容
    text = md_path.read_text(encoding="utf-8")

    # 删除 markdown 图片
    text = re.sub(r'!\[.*?\]\(.*?\)', '', text)

    # 删除 html 注释
    text = re.sub(r'<!--.*?-->', '', text)

    # 清理多余空行
    text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)

    md_path.write_text(text, encoding="utf-8")

    # 3. 删除所有图片
    for img in md_dir.glob("*.jpeg"):
        img.unlink()

    for img in md_dir.glob("*.png"):
        img.unlink()

    # 4. 删除 meta json
    meta_file = md_dir / f"{pdf_path.stem}_meta.json"
    if meta_file.exists():
        meta_file.unlink()

    # 5. 把 md 移动到外层目录
    final_md = output_dir / f"{pdf_path.stem}.md"
    md_path.rename(final_md)

    # 6. 删除 marker 生成的目录
    md_dir.rmdir()

    return final_md


if __name__ == "__main__":
    pdf_path = "/home/guoziyang/AIgorithm_Agent/src/preparation/parser/tmp/1_input.pdf"
    md_path = pdf_to_md(pdf_path)
    ppt_path = "/home/guoziyang/AIgorithm_Agent/src/preparation/parser/tmp/中期汇报.pptx"
    md_path2 = ppt_to_md(ppt_path)
    print(f"Markdown file generated at: {md_path}")
    print(f"Markdown file generated at: {md_path2}")